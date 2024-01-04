import json
import re

import docx2txt
import pdfminer.high_level
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_pptx_text(filename: str) -> str:
    def check_recursively_for_text(this_set_of_shapes, text_run):
        for shape in this_set_of_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                check_recursively_for_text(shape.shapes, text_run)
            else:
                if hasattr(shape, "text"):
                    text_run.append(shape.text)
        return text_run

    output = {}
    presentation = pptx.Presentation(filename)
    for i, slide in enumerate(presentation.slides):
        slide_number = i + 1
        text_run = []
        text_run = check_recursively_for_text(slide.shapes, text_run)
        output[f'Slide {slide_number}'] = re.sub('[　 \\t]+', ' ', ' '.join(text_run).strip())
    return json.dumps(output, ensure_ascii=False)


def read_pdf_text(filename: str) -> str:
    return json.dumps(
        re.sub('[　 \\t]+', ' ', re.sub("\\n+", '\n', pdfminer.high_level.extract_text(filename))).strip(),
        ensure_ascii=False)


def read_docx_text(filename: str) -> str:
    return json.dumps(
        re.sub('[　 \\t]+', ' ', re.sub("\\n+", '\n', docx2txt.process(filename))).strip(),
        ensure_ascii=False)
